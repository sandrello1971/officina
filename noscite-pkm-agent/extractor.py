"""
Estrazione testo e metadati tecnici da qualsiasi tipo di file.
"""

import logging
import subprocess
from pathlib import Path

log = logging.getLogger(__name__)


def extract(path: Path) -> dict:
    """
    Restituisce un dizionario con:
    - text: testo estratto (può essere vuoto per media)
    - tech_meta: metadati tecnici (dimensione, durata, ecc.)
    - extraction_method: metodo usato
    """
    ext = path.suffix.lower()
    size_bytes = path.stat().st_size
    size_str = _human_size(size_bytes)

    result = {
        "text": "",
        "tech_meta": {"file_size": size_str, "extension": ext},
        "extraction_method": "none",
    }

    try:
        if ext == ".pdf":
            result.update(_extract_pdf(path))
        elif ext == ".docx":
            result.update(_extract_docx(path))
        elif ext == ".pptx":
            result.update(_extract_pptx(path))
        elif ext in (".xlsx", ".xls"):
            result.update(_extract_excel(path))
        elif ext in (".txt", ".csv", ".rtf", ".md"):
            result.update(_extract_plain(path))
        elif ext in (".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm"):
            result.update(_extract_video_meta(path))
        elif ext in (".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a", ".wma"):
            result.update(_extract_audio_meta(path))
        elif ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"):
            result.update(_extract_image_meta(path))
        else:
            result["extraction_method"] = "unsupported"
    except Exception as e:
        log.warning(f"Estrazione fallita per {path.name}: {e}")
        result["extraction_error"] = str(e)

    return result


# ── PDF ──────────────────────────────────────

def _extract_pdf(path: Path) -> dict:
    try:
        import pdfplumber
        pages_text = []
        with pdfplumber.open(path) as pdf:
            num_pages = len(pdf.pages)
            for page in pdf.pages[:30]:  # max 30 pagine
                t = page.extract_text()
                if t:
                    pages_text.append(t)
        return {
            "text": "\n".join(pages_text),
            "tech_meta": {"pages": num_pages},
            "extraction_method": "pdfplumber",
        }
    except ImportError:
        pass

    # Fallback: pdfminer
    try:
        from pdfminer.high_level import extract_text
        text = extract_text(str(path))
        return {"text": text or "", "extraction_method": "pdfminer"}
    except ImportError:
        return {"text": "", "extraction_method": "pdf_unavailable"}


# ── DOCX ─────────────────────────────────────

def _extract_docx(path: Path) -> dict:
    try:
        import docx
        doc = docx.Document(str(path))

        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())

        tables_md = []
        for ti, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                cells = [c.text.strip().replace("\n", " ") or " " for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, sep)
                tables_md.append(f"\n**Tabella {ti}:**\n\n" + "\n".join(rows))

        images = []
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref.lower():
                    images.append(rel.target_ref.split("/")[-1])
        except Exception:
            pass

        text = "\n\n".join(parts)
        if tables_md:
            text += "\n\n" + "\n\n".join(tables_md)
        if images:
            text += "\n\n**Immagini incorporate:** " + ", ".join(f"`{img}`" for img in images)

        return {
            "text": text,
            "tech_meta": {
                "paragraphs": len([p for p in doc.paragraphs if p.text.strip()]),
                "tables": len(doc.tables),
                "images": len(images),
            },
            "extraction_method": "python-docx",
        }
    except ImportError:
        return {"text": "", "extraction_method": "docx_unavailable"}


# ── PPTX ─────────────────────────────────────

def _extract_pptx(path: Path) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text.strip())
        return {
            "text": "\n".join(slides_text),
            "tech_meta": {"slides": len(prs.slides)},
            "extraction_method": "python-pptx",
        }
    except ImportError:
        return {"text": "", "extraction_method": "pptx_unavailable"}


# ── EXCEL ────────────────────────────────────

def _extract_excel(path: Path) -> dict:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets_data = []
        for sheet in wb.sheetnames[:5]:  # max 5 fogli
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(max_row=50, values_only=True):
                row_str = " | ".join(str(c) for c in row if c is not None)
                if row_str.strip():
                    rows.append(row_str)
            sheets_data.append(f"[Foglio: {sheet}]\n" + "\n".join(rows))
        return {
            "text": "\n\n".join(sheets_data),
            "tech_meta": {"sheets": len(wb.sheetnames)},
            "extraction_method": "openpyxl",
        }
    except ImportError:
        return {"text": "", "extraction_method": "excel_unavailable"}


# ── PLAIN TEXT ───────────────────────────────

def _extract_plain(path: Path) -> dict:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return {"text": text, "extraction_method": "plain_text"}
    except Exception as e:
        return {"text": "", "extraction_method": "plain_failed", "error": str(e)}


# ── VIDEO ────────────────────────────────────

def _extract_video_meta(path: Path) -> dict:
    meta = _ffprobe_meta(path)
    return {"text": "", "tech_meta": meta, "extraction_method": "ffprobe"}


# ── AUDIO ────────────────────────────────────

def _extract_audio_meta(path: Path) -> dict:
    meta = _ffprobe_meta(path)
    return {"text": "", "tech_meta": meta, "extraction_method": "ffprobe"}


# ── IMMAGINE ─────────────────────────────────

def _extract_image_meta(path: Path) -> dict:
    meta = {}
    try:
        from PIL import Image
        with Image.open(path) as img:
            meta["width"] = img.width
            meta["height"] = img.height
            meta["mode"] = img.mode
            meta["format"] = img.format
        return {"text": "", "tech_meta": meta, "extraction_method": "Pillow"}
    except ImportError:
        return {"text": "", "tech_meta": meta, "extraction_method": "image_unavailable"}


# ── UTILITY ──────────────────────────────────

def _ffprobe_meta(path: Path) -> dict:
    """Usa ffprobe per estrarre metadati da file media."""
    try:
        import json
        result = subprocess.run(
            [
                "ffprobe", "-v", "quiet", "-print_format", "json",
                "-show_format", "-show_streams", str(path),
            ],
            capture_output=True, text=True, timeout=10,
        )
        if result.returncode == 0:
            data = json.loads(result.stdout)
            fmt = data.get("format", {})
            duration = float(fmt.get("duration", 0))
            return {
                "duration_seconds": round(duration),
                "duration_hms": _seconds_to_hms(duration),
                "bit_rate": fmt.get("bit_rate", "N/A"),
                "format_name": fmt.get("format_long_name", ""),
            }
    except (FileNotFoundError, Exception):
        pass
    return {}


def _seconds_to_hms(seconds: float) -> str:
    s = int(seconds)
    h, remainder = divmod(s, 3600)
    m, sec = divmod(remainder, 60)
    return f"{h:02d}:{m:02d}:{sec:02d}"


def _human_size(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if n < 1024:
            return f"{n:.1f} {unit}"
        n /= 1024
    return f"{n:.1f} TB"
