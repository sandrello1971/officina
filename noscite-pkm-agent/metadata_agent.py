#!/usr/bin/env python3
"""
Obsidian AI Metadata Agent
==========================
Monitora /vault/originals/, estrae metadati con Claude AI,
genera file .md in /vault/metadata/ con YAML frontmatter,
tag semantici e wikilink per il grafo di Obsidian.

Uso:
    python metadata_agent.py --watch          # monitoraggio continuo
    python metadata_agent.py --scan           # scansione una tantum
    python metadata_agent.py --file doc.pdf   # singolo file
    python metadata_agent.py --rebuild-index  # rigenera tabella ricerca
"""

import os
import json
import time
import hashlib
import argparse
import re
from pathlib import Path
from datetime import datetime

# ---------------------------------------------------------------------------
# Configurazione — modifica questi percorsi per il tuo vault Obsidian
# ---------------------------------------------------------------------------
CONFIG = {
    "originals_dir": "./vault/originals",       # cartella file originali
    "metadata_dir": "./vault/metadata",          # cartella .md generati
    "index_file": "./vault/metadata/_INDEX.md",  # tabella ricerca globale
    "processed_db": "./vault/.processed.json",   # DB file già elaborati
    "api_key_env": "ANTHROPIC_API_KEY",          # variabile env per la chiave API
    "model": "claude-sonnet-4-20250514",
    "max_text_chars": 8000,                      # max caratteri testo estratto
    "watch_interval": 5,                         # secondi tra scansioni in watch mode
}

# Estensioni supportate e come estrarre testo
SUPPORTED_EXTENSIONS = {
    # Documenti
    ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".xlsx": "xlsx", ".xls": "xlsx", ".csv": "csv",
    ".txt": "text", ".md": "text", ".rtf": "text",
    # Immagini
    ".jpg": "image", ".jpeg": "image", ".png": "image",
    ".gif": "image", ".webp": "image", ".bmp": "image",
    # Audio
    ".mp3": "audio", ".m4a": "audio", ".wav": "audio",
    ".flac": "audio", ".ogg": "audio", ".aac": "audio",
    # Video
    ".mp4": "video", ".mov": "video", ".avi": "video",
    ".mkv": "video", ".wmv": "video", ".webm": "video",
    # Archivi e codice
    ".zip": "archive", ".tar": "archive", ".gz": "archive",
    ".py": "code", ".js": "code", ".ts": "code",
    ".html": "code", ".css": "code", ".json": "code",
    ".xml": "code", ".yaml": "code", ".yml": "code",
}


# ---------------------------------------------------------------------------
# Estrazione testo per tipo di file
# ---------------------------------------------------------------------------

def extract_text_pdf(path: Path) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages[:20]):  # max 20 pagine
            text = page.extract_text()
            if text:
                pages.append(f"[Pagina {i+1}]\n{text}")
        return "\n\n".join(pages)
    except Exception as e:
        return f"[Errore estrazione PDF: {e}]"


def extract_text_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Includi anche testo dalle tabelle
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)
    except Exception as e:
        return f"[Errore estrazione DOCX: {e}]"


def extract_text_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[Errore estrazione PPTX: {e}]"


def extract_text_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames[:5]:  # max 5 fogli
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=50, values_only=True):  # max 50 righe
                row_vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in row_vals):
                    rows.append(" | ".join(row_vals))
            if rows:
                sheets.append(f"[Foglio: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    except Exception as e:
        return f"[Errore estrazione XLSX: {e}]"


def extract_text_csv(path: Path) -> str:
    try:
        import csv
        rows = []
        with open(path, encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i > 100:
                    rows.append(f"... ({i} righe totali)")
                    break
                rows.append(" | ".join(row))
        return "\n".join(rows)
    except Exception as e:
        return f"[Errore estrazione CSV: {e}]"


def extract_text_plain(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[Errore lettura file: {e}]"


def get_image_info(path: Path) -> str:
    try:
        from PIL import Image
        img = Image.open(path)
        return f"Immagine {img.format}, dimensioni: {img.size[0]}x{img.size[1]}px, modalità: {img.mode}"
    except Exception as e:
        return f"File immagine (dettagli non disponibili: {e})"


def get_media_info(path: Path, media_type: str) -> str:
    """Metadati base per audio/video tramite mutagen se disponibile."""
    info = f"File {media_type.upper()}, dimensione: {path.stat().st_size / 1024 / 1024:.1f} MB"
    try:
        import mutagen
        audio = mutagen.File(str(path))
        if audio and audio.tags:
            tags = {k: str(v) for k, v in list(audio.tags.items())[:10]}
            info += f"\nTag: {json.dumps(tags, ensure_ascii=False)}"
        if hasattr(audio, 'info') and audio.info:
            if hasattr(audio.info, 'length'):
                mins = int(audio.info.length // 60)
                secs = int(audio.info.length % 60)
                info += f"\nDurata: {mins}:{secs:02d}"
    except Exception:
        pass
    return info


def extract_content(path: Path) -> tuple[str, str]:
    """
    Estrae contenuto testuale e tipo del file.
    Ritorna (contenuto, tipo_file).
    """
    ext = path.suffix.lower()
    file_type = SUPPORTED_EXTENSIONS.get(ext, "unknown")

    if file_type == "pdf":
        content = extract_text_pdf(path)
    elif file_type == "docx":
        content = extract_text_docx(path)
    elif file_type == "pptx":
        content = extract_text_pptx(path)
    elif file_type == "xlsx":
        content = extract_text_xlsx(path)
    elif file_type == "csv":
        content = extract_text_csv(path)
    elif file_type == "text":
        content = extract_text_plain(path)
    elif file_type == "code":
        content = extract_text_plain(path)
    elif file_type == "image":
        content = get_image_info(path)
    elif file_type in ("audio", "video"):
        content = get_media_info(path, file_type)
    elif file_type == "archive":
        try:
            import zipfile
            if zipfile.is_zipfile(path):
                with zipfile.ZipFile(path) as zf:
                    content = "Contenuto archivio ZIP:\n" + "\n".join(zf.namelist()[:50])
            else:
                content = "File archivio compresso"
        except Exception:
            content = "File archivio"
    else:
        content = f"File di tipo sconosciuto ({ext})"

    # Limita lunghezza
    if len(content) > CONFIG["max_text_chars"]:
        content = content[:CONFIG["max_text_chars"]] + "\n\n[... contenuto troncato ...]"

    return content, file_type


# ---------------------------------------------------------------------------
# Claude AI: generazione metadati
# ---------------------------------------------------------------------------

def call_claude_api(content_text: str, filename: str, file_type: str, file_size_kb: float) -> dict:
    """
    Chiama l'API di Claude per generare metadati strutturati.
    Richiede la variabile d'ambiente ANTHROPIC_API_KEY.
    """
    import urllib.request
    import urllib.error

    api_key = os.environ.get(CONFIG["api_key_env"])
    if not api_key:
        raise ValueError(
            f"Variabile d'ambiente {CONFIG['api_key_env']} non trovata.\n"
            f"Esporta la chiave: export ANTHROPIC_API_KEY='sk-ant-...'"
        )

    system_prompt = """Sei un assistente specializzato nell'estrazione di metadati da documenti per un sistema di gestione della conoscenza personale (PKM) basato su Obsidian.

Analizza il contenuto fornito e restituisci SOLO un oggetto JSON valido (senza markdown, senza backtick) con questa struttura esatta:

{
  "titolo": "titolo descrittivo del documento (non il nome file)",
  "sommario": "riassunto in 2-4 frasi che cattura l'essenza del documento",
  "tipo_documento": "uno tra: report, presentazione, foglio_dati, immagine, audio, video, codice, documento, articolo, email, contratto, manuale, nota, visura, verbale, certificato, vademecum, agenda, faq, viaggio, spettacolo, schema, progetto, altro",
  "argomento_principale": "argomento primario del documento in 3-5 parole",
  "argomenti_secondari": ["lista", "di", "argomenti", "secondari"],
  "lingua": "it/en/fr/de/es/...",
  "data_stimata": "anno stimato di creazione (YYYY) o null se non determinabile",
  "autori": ["autore1", "autore2"] ,
  "organizzazione": "azienda/ente/istituzione se presente, altrimenti null",
  "parole_chiave": ["keyword1", "keyword2", "keyword3", "keyword4", "keyword5"],
  "tag_obsidian": ["#tag1", "#tag2", "#tag3", "#tag4", "#tag5", "#tag6"],
  "livello_importanza": "alta/media/bassa (stima soggettiva basata sul contenuto)",
  "stato": "completo/bozza/frammento",
  "note_agente": "osservazioni specifiche sull'analisi di questo documento"
}

Per i tag Obsidian, usa tag gerarchici dove appropriato (es. #finanza/bilancio, #progetto/alfa, #anno/2024).
Genera almeno 5-8 tag utili per creare relazioni nel grafo di Obsidian."""

    user_message = f"""Analizza questo file e genera i metadati richiesti.

Nome file: {filename}
Tipo: {file_type}
Dimensione: {file_size_kb:.1f} KB

Contenuto estratto:
---
{content_text}
---

Restituisci SOLO il JSON, senza alcun testo aggiuntivo."""

    payload = json.dumps({
        "model": CONFIG["model"],
        "max_tokens": 1500,
        "system": system_prompt,
        "messages": [{"role": "user", "content": user_message}]
    }).encode("utf-8")

    req = urllib.request.Request(
        "https://api.anthropic.com/v1/messages",
        data=payload,
        headers={
            "Content-Type": "application/json",
            "x-api-key": api_key,
            "anthropic-version": "2023-06-01",
        },
        method="POST"
    )

    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read().decode("utf-8"))
            raw = data["content"][0]["text"].strip()
            # Pulisci eventuale markdown residuo
            raw = re.sub(r'^```json\s*', '', raw)
            raw = re.sub(r'\s*```$', '', raw)
            return json.loads(raw)
    except urllib.error.HTTPError as e:
        error_body = e.read().decode("utf-8")
        raise RuntimeError(f"Errore API Claude {e.code}: {error_body}")
    except json.JSONDecodeError as e:
        raise RuntimeError(f"Risposta JSON non valida da Claude: {e}\nRisposta: {raw[:200]}")


# ---------------------------------------------------------------------------
# Generazione file .md per Obsidian
# ---------------------------------------------------------------------------

def generate_metadata_md(
    source_path: Path,
    metadata: dict,
    content_preview: str,
    file_type: str,
    processing_date: str,
) -> str:
    """
    Genera il contenuto del file .md con YAML frontmatter, 
    corpo ricco e wikilink per Obsidian.
    """
    # Prepara tag per YAML (rimuovi # iniziale per il frontmatter)
    tags_yaml = metadata.get("tag_obsidian", [])
    tags_clean = [t.lstrip("#") for t in tags_yaml]

    # Aggiunge tag base tipo e lingua
    base_tags = [
        f"tipo/{file_type}",
        f"lingua/{metadata.get('lingua', 'unknown')}",
        "catalogo/documento"
    ]
    if metadata.get("data_stimata"):
        base_tags.append(f"anno/{metadata['data_stimata']}")
    all_tags = base_tags + tags_clean
    tags_yaml_str = "\n".join(f"  - {t}" for t in all_tags)

    # Parole chiave per aliases
    keywords = metadata.get("parole_chiave", [])
    aliases_str = "\n".join(f'  - "{k}"' for k in keywords[:3])

    # Percorso relativo per il link al file originale
    orig_rel = source_path.name

    # Autori
    autori = metadata.get("autori", [])
    autori_str = ", ".join(autori) if autori else "non specificato"

    # Argomenti secondari come wikilink
    argomenti_links = " | ".join(
        f"[[{a}]]" for a in metadata.get("argomenti_secondari", [])[:5]
    )

    # Frontmatter YAML
    frontmatter = f"""---
title: "{metadata.get('titolo', source_path.stem)}"
aliases:
{aliases_str}
tags:
{tags_yaml_str}
tipo_documento: {metadata.get('tipo_documento', 'altro')}
lingua: {metadata.get('lingua', 'unknown')}
argomento_principale: "{metadata.get('argomento_principale', '')}"
data_stimata: {metadata.get('data_stimata', 'null')}
autori: [{', '.join(f'"{a}"' for a in autori)}]
organizzazione: "{metadata.get('organizzazione', '') or ''}"
livello_importanza: {metadata.get('livello_importanza', 'media')}
stato_documento: {metadata.get('stato', 'completo')}
file_originale: "originals/{orig_rel}"
dimensione_kb: {source_path.stat().st_size / 1024:.1f}
data_catalogazione: {processing_date}
agente_versione: "1.0"
---"""

    # Corpo del documento
    body = f"""
# {metadata.get('titolo', source_path.stem)}

> {metadata.get('sommario', '')}

---

## 📋 Informazioni documento

| Campo | Valore |
|-------|--------|
| **Tipo** | {metadata.get('tipo_documento', file_type)} |
| **Argomento** | {metadata.get('argomento_principale', '')} |
| **Autori** | {autori_str} |
| **Organizzazione** | {metadata.get('organizzazione', 'n/d') or 'n/d'} |
| **Anno stimato** | {metadata.get('data_stimata', 'n/d') or 'n/d'} |
| **Lingua** | {metadata.get('lingua', 'n/d')} |
| **Importanza** | {metadata.get('livello_importanza', 'media')} |
| **Stato** | {metadata.get('stato', 'completo')} |
| **File originale** | [[originals/{orig_rel}]] |

---

## 🗝️ Parole chiave

{" • ".join(f"`{k}`" for k in keywords)}

---

## 🔗 Argomenti correlati

{argomenti_links if argomenti_links else "_Nessun argomento secondario identificato_"}

---

## 📝 Note agente AI

_{metadata.get('note_agente', '')}_

---

## 👁️ Anteprima contenuto

```
{content_preview[:500].strip()}{"..." if len(content_preview) > 500 else ""}
```

---

*Catalogato automaticamente il {processing_date} · [[_INDEX|← Torna all'indice]]*
"""

    return frontmatter + body


# ---------------------------------------------------------------------------
# Gestione database file processati
# ---------------------------------------------------------------------------

class ProcessedDB:
    def __init__(self, db_path: str):
        self.path = Path(db_path)
        self.data: dict = {}
        self.load()

    def load(self):
        if self.path.exists():
            with open(self.path, encoding="utf-8") as f:
                self.data = json.load(f)

    def save(self):
        self.path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.path, "w", encoding="utf-8") as f:
            json.dump(self.data, f, ensure_ascii=False, indent=2)

    def file_hash(self, path: Path) -> str:
        h = hashlib.md5()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()

    def is_processed(self, path: Path) -> bool:
        key = str(path)
        if key not in self.data:
            return False
        stored_hash = self.data[key].get("hash")
        current_hash = self.file_hash(path)
        return stored_hash == current_hash

    def mark_processed(self, path: Path, metadata_file: str):
        self.data[str(path)] = {
            "hash": self.file_hash(path),
            "metadata_file": metadata_file,
            "processed_at": datetime.now().isoformat(),
        }
        self.save()

    def all_entries(self) -> list:
        return list(self.data.values())


# ---------------------------------------------------------------------------
# Generazione indice globale
# ---------------------------------------------------------------------------

def rebuild_index(metadata_dir: Path, index_path: Path):
    """Scansiona tutti i .md e rigenera la tabella _INDEX.md."""

    entries = []
    for md_file in sorted(metadata_dir.glob("*.md")):
        if md_file.name.startswith("_"):
            continue
        content = md_file.read_text(encoding="utf-8")
        # Estrai campi da frontmatter
        def get_field(field, default=""):
            m = re.search(rf'^{field}:\s*"?([^"\n]+)"?', content, re.MULTILINE)
            return m.group(1).strip() if m else default

        entries.append({
            "link": f"[[{md_file.stem}]]",
            "titolo": get_field("title"),
            "tipo": get_field("tipo_documento"),
            "argomento": get_field("argomento_principale"),
            "anno": get_field("data_stimata", "n/d"),
            "lingua": get_field("lingua"),
            "importanza": get_field("livello_importanza"),
        })

    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    rows = "\n".join(
        f"| {e['link']} | {e['tipo']} | {e['argomento']} | {e['anno']} | {e['lingua']} | {e['importanza']} |"
        for e in entries
    )

    index_content = f"""---
tags:
  - catalogo/indice
  - sistema/pkm
title: "Indice documenti"
---

# 📚 Indice documenti PKM

> Tabella generata automaticamente — {len(entries)} documenti catalogati  
> Ultimo aggiornamento: {now}

---

## 🔍 Ricerca veloce

Usa il plugin **Dataview** per query avanzate:

```dataview
TABLE tipo_documento AS Tipo, argomento_principale AS Argomento, data_stimata AS Anno
FROM "metadata"
WHERE livello_importanza = "alta"
SORT data_stimata DESC
```

---

## 📋 Tutti i documenti

| Documento | Tipo | Argomento | Anno | Lingua | Importanza |
|-----------|------|-----------|------|--------|------------|
{rows}

---

## 🏷️ Tag cloud

Usa **Tag Pane** nella sidebar di Obsidian per navigare per tag.  
Usa **Graph View** per visualizzare le relazioni tra documenti.

---

*Aggiornato automaticamente dall'AI Metadata Agent*
"""

    index_path.parent.mkdir(parents=True, exist_ok=True)
    index_path.write_text(index_content, encoding="utf-8")
    print(f"  ✅ Indice aggiornato: {len(entries)} documenti → {index_path}")


# ---------------------------------------------------------------------------
# Elaborazione di un singolo file
# ---------------------------------------------------------------------------

def process_file(source_path: Path, db: ProcessedDB, verbose: bool = True) -> bool:
    """
    Elabora un file: estrae contenuto, genera metadati AI, salva .md.
    Ritorna True se il file è stato processato, False se già aggiornato.
    """
    if db.is_processed(source_path):
        if verbose:
            print(f"  ⏭️  Già catalogato: {source_path.name}")
        return False

    print(f"  🔍 Elaborazione: {source_path.name}")

    ext = source_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        print(f"  ⚠️  Estensione non supportata: {ext} — file ignorato")
        return False

    # 1. Estrai contenuto
    try:
        content, file_type = extract_content(source_path)
    except Exception as e:
        print(f"  ❌ Errore estrazione contenuto: {e}")
        return False

    # 2. Chiama Claude AI
    try:
        file_size_kb = source_path.stat().st_size / 1024
        metadata = call_claude_api(content, source_path.name, file_type, file_size_kb)
    except Exception as e:
        print(f"  ❌ Errore API Claude: {e}")
        return False

    # 3. Genera file .md
    processing_date = datetime.now().strftime("%Y-%m-%d")
    md_content = generate_metadata_md(
        source_path, metadata, content, file_type, processing_date
    )

    # Nome file .md basato sul nome originale (safe)
    safe_name = re.sub(r'[^\w\-\. ]', '_', source_path.stem)
    safe_name = re.sub(r'\s+', '_', safe_name)
    md_filename = f"{safe_name}.md"

    metadata_dir = Path(CONFIG["metadata_dir"])
    metadata_dir.mkdir(parents=True, exist_ok=True)
    md_path = metadata_dir / md_filename

    # Evita sovrascrittura di file con lo stesso nome da file diversi
    if md_path.exists():
        counter = 1
        while md_path.exists():
            md_filename = f"{safe_name}_{counter}.md"
            md_path = metadata_dir / md_filename
            counter += 1

    md_path.write_text(md_content, encoding="utf-8")

    # 4. Aggiorna database
    db.mark_processed(source_path, str(md_path))

    titolo = metadata.get('titolo', source_path.stem)[:50]
    tags_str = " ".join(metadata.get("tag_obsidian", [])[:4])
    print(f"  ✅ Catalogato: {titolo}")
    print(f"     Tag: {tags_str}")
    print(f"     → {md_path}")

    return True


# ---------------------------------------------------------------------------
# Modalità di esecuzione
# ---------------------------------------------------------------------------

def scan_once(verbose: bool = True) -> int:
    """Scansione singola di tutti i file in originals_dir."""
    originals_dir = Path(CONFIG["originals_dir"])
    if not originals_dir.exists():
        print(f"❌ Cartella originals non trovata: {originals_dir}")
        print("   Crea la cartella e aggiungi i tuoi documenti.")
        return 0

    db = ProcessedDB(CONFIG["processed_db"])
    processed = 0

    files = [f for f in originals_dir.iterdir() if f.is_file()]
    total = len(files)
    print(f"\n📂 Scansione {originals_dir} — {total} file trovati\n")

    for i, file_path in enumerate(sorted(files), 1):
        print(f"[{i}/{total}] ", end="")
        if process_file(file_path, db, verbose):
            processed += 1

    if processed > 0:
        print("\n🔄 Rigenero indice...")
        rebuild_index(Path(CONFIG["metadata_dir"]), Path(CONFIG["index_file"]))

    print(f"\n✨ Completato: {processed}/{total} file catalogati")
    return processed


def watch_mode():
    """Monitoraggio continuo della cartella originals."""
    print(f"👁️  Modalità watch avviata (controllo ogni {CONFIG['watch_interval']}s)")
    print(f"   Monitorando: {CONFIG['originals_dir']}")
    print("   Premi Ctrl+C per fermare\n")

    try:
        while True:
            scan_once(verbose=False)
            time.sleep(CONFIG["watch_interval"])
    except KeyboardInterrupt:
        print("\n\n⏹️  Watch fermato.")


def process_single_file(file_path: str):
    """Elabora un singolo file specificato."""
    path = Path(file_path)
    if not path.exists():
        print(f"❌ File non trovato: {file_path}")
        return
    db = ProcessedDB(CONFIG["processed_db"])
    # Forza rielaborazione
    if str(path) in db.data:
        del db.data[str(path)]
    process_file(path, db)
    rebuild_index(Path(CONFIG["metadata_dir"]), Path(CONFIG["index_file"]))


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Obsidian AI Metadata Agent — cataloga documenti con Claude AI"
    )
    parser.add_argument("--watch", action="store_true", help="Monitoraggio continuo")
    parser.add_argument("--scan", action="store_true", help="Scansione singola")
    parser.add_argument("--file", type=str, help="Elabora un singolo file")
    parser.add_argument("--rebuild-index", action="store_true", help="Rigenera solo l'indice")
    parser.add_argument("--config", type=str, help="Mostra configurazione attuale")
    args = parser.parse_args()

    print("=" * 60)
    print("  🤖 Obsidian AI Metadata Agent")
    print("=" * 60)

    if args.config or (not any([args.watch, args.scan, args.file, args.rebuild_index])):
        print("\nConfigurazione attuale:")
        for k, v in CONFIG.items():
            print(f"  {k}: {v}")
        print("\nUso:")
        print("  python metadata_agent.py --scan           # scansione una tantum")
        print("  python metadata_agent.py --watch          # monitoraggio continuo")
        print("  python metadata_agent.py --file doc.pdf   # singolo file")
        print("  python metadata_agent.py --rebuild-index  # rigenera indice")
        print("\nVariabile richiesta: export ANTHROPIC_API_KEY='sk-ant-...'")
        return

    if args.rebuild_index:
        rebuild_index(Path(CONFIG["metadata_dir"]), Path(CONFIG["index_file"]))
    elif args.file:
        process_single_file(args.file)
    elif args.watch:
        watch_mode()
    elif args.scan:
        scan_once()


if __name__ == "__main__":
    main()
