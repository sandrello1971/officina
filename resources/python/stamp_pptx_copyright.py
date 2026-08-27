#!/usr/bin/env python
"""Stampiglia la dicitura di copyright in fondo a OGNI slide di un .pptx ESISTENTE.

Gemello di build_pptx.py::footer, ma per le presentazioni CARICATE a mano
dall'admin: quelle generate ricevono il footer in fase di render, queste no, e
finirebbero davanti agli studenti (e dentro i video, che nascono dai PNG delle
slide) senza tutela.

Uso:  stamp_pptx_copyright.py <file.pptx> <notice>
Il file viene riscritto in place. Exit 0 = stampigliato, exit != 0 = nulla di
fatto (il chiamante PHP tratta il fallimento come bloccante: contenuto senza
dicitura non si pubblica).

Idempotente: le slide già stampigliate sono riconosciute dal marcatore nel nome
della shape, così un ri-upload o un re-stamp non accumula scritte sovrapposte.
"""
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Nome della shape: marcatore di idempotenza (le shape di python-pptx lo espongono).
SHAPE_NAME = "nosc-copyright"
HEIGHT = Inches(0.32)
MARGIN = Inches(0.55)
# Grigio medio: leggibile sia su fondo chiaro sia su fondo scuro, senza dover
# indovinare lo sfondo di una slide che non abbiamo costruito noi.
COLOR = RGBColor(0x8A, 0x8A, 0x8A)


def already_stamped(slide):
    return any(getattr(shape, "name", "") == SHAPE_NAME for shape in slide.shapes)


def stamp(slide, prs, notice):
    width = prs.slide_width - (2 * MARGIN)
    top = prs.slide_height - HEIGHT - Inches(0.10)
    box = slide.shapes.add_textbox(MARGIN, Emu(int(top)), Emu(int(width)), HEIGHT)
    box.name = SHAPE_NAME

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = notice
    run.font.size = Pt(7)
    run.font.color.rgb = COLOR


def main():
    if len(sys.argv) < 3:
        print("uso: stamp_pptx_copyright.py <file.pptx> <notice>", file=sys.stderr)
        return 2

    path = sys.argv[1]
    notice = sys.argv[2].strip()
    if not notice:
        print("dicitura vuota: niente da stampigliare", file=sys.stderr)
        return 3

    prs = Presentation(path)
    for slide in prs.slides:
        if not already_stamped(slide):
            stamp(slide, prs, notice)
    prs.save(path)
    return 0


if __name__ == "__main__":
    sys.exit(main())
