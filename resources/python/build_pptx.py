#!/usr/bin/env python3
"""
Genera un .pptx da una specifica JSON (Schola P21). Non reinventa il formato:
usa python-pptx (OOXML). La spec arriva da stdin, il path di output da argv[1].

Spec JSON:
{
  "title": "Titolo lezione",
  "subtitle": "Argomento · Materia",
  "school": "Nome scuola o piattaforma",     # branding (titolo slide iniziale)
  "accent": "55B1AE",                          # colore accento esadecimale (no '#')
  "slides": [
    {"title": "Sezione", "bullets": ["punto 1", "punto 2", ...], "notes": "..."},
    ...
  ]
}
"""
import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def hex_to_rgb(value, fallback=(85, 177, 174)):
    try:
        value = (value or "").lstrip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except Exception:
        return RGBColor(*fallback)


def main():
    spec = json.load(sys.stdin)
    out_path = sys.argv[1]

    accent = hex_to_rgb(spec.get("accent"))
    prs = Presentation()

    # --- Slide di titolo (branding) ---
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = spec.get("title", "Lezione")
    subtitle = slide.placeholders[1]
    parts = [p for p in [spec.get("subtitle"), spec.get("school")] if p]
    subtitle.text = "\n".join(parts)
    try:
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent
    except Exception:
        pass

    # --- Una slide per sezione: titolo + bullet ---
    bullet_layout = prs.slide_layouts[1]  # Title and Content
    for s in spec.get("slides", []):
        sl = prs.slides.add_slide(bullet_layout)
        sl.shapes.title.text = s.get("title", "")
        try:
            sl.shapes.title.text_frame.paragraphs[0].font.color.rgb = accent
        except Exception:
            pass

        body = sl.placeholders[1].text_frame
        body.clear()
        bullets = s.get("bullets") or []
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(b)
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)

        notes = s.get("notes")
        if notes:
            sl.notes_slide.notes_text_frame.text = str(notes)

    prs.save(out_path)
    print(json.dumps({"ok": True, "slides": len(spec.get("slides", [])) + 1}))


if __name__ == "__main__":
    main()
